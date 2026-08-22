from __future__ import annotations

import json
import io
import hashlib
import os
import re
import stat
import unicodedata
from dataclasses import asdict, dataclass, field, replace
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup

from app.core.config import get_settings
from app.core.utils import source_type_from_path
from app.services.storage import (
    FROZEN_SOURCE_SNAPSHOT_PROTOCOL_VERSION,
    FrozenSourceSnapshot,
    VerifiedSourceIdentity,
    _TEST_DURABILITY_ADAPTER,
    _validate_seekable_upload_content,
    validate_frozen_source_snapshot,
)


PARSER_LAYOUT_PROTOCOL_VERSION = "prepared_document_layout_v1"
FLOW_BLOCK_PROTOCOL_VERSION = "deterministic_flow_blocks_v14"
MARKDOWN_PIPE_TABLE_PROTOCOL_VERSION = "markdown_pipe_table_v5"
MARKDOWN_BLOCK_START_PROTOCOL_VERSION = "markdown_block_start_precedence_v11"
MARKDOWN_TABLE_ROW_NORMALIZATION_PROTOCOL_VERSION = (
    "markdown_table_rows_pad_missing_truncate_excess_v2"
)
MARKDOWN_INDENTED_CODE_PROTOCOL_VERSION = "markdown_indented_code_v1"
MARKDOWN_HTML_BLOCK_PROTOCOL_VERSION = "markdown_html_blocks_gfm_0_29_v2"
MARKDOWN_LINK_REFERENCE_PROTOCOL_VERSION = "markdown_link_reference_gfm_0_29_v2"
MARKDOWN_TARGET_SPEC_VERSION = "gfm_0_29"

_MARKDOWN_TABLE_DELIMITER_CELL_RE = re.compile(r"^:?-{3,}:?$")
_MARKDOWN_FENCE_OPEN_RE = re.compile(r"^(`{3,}|~{3,})(.*)$")
_MARKDOWN_ATX_OPEN_RE = re.compile(
    r"^(?P<indent> {0,3})(?P<opening>#+)(?P<remainder>.*)$"
)
_MARKDOWN_ATX_CLOSING_RE = re.compile(r"[ \t]+#+[ \t]*$")
_MARKDOWN_LIST_MARKER_RE = re.compile(
    r"^(?:(?P<bullet>[*+-])|(?P<number>\d{1,9})(?P<ordered_delimiter>[.)]))"
    r"(?P<spacing>[ \t]*)(?P<content>.*)$"
)
_MARKDOWN_HTML_TYPE_1_RE = re.compile(
    r"^<(?:script|pre|style)(?:[ \t]+|>|$)",
    flags=re.IGNORECASE,
)
_MARKDOWN_HTML_TYPE_1_END_RE = re.compile(
    r"</(?:script|pre|style)>",
    flags=re.IGNORECASE,
)
_MARKDOWN_HTML_BLOCK_TAG_RE = re.compile(
    r"^</?(?:address|article|aside|base|basefont|blockquote|body|caption|center|"
    r"col|colgroup|dd|details|dialog|dir|div|dl|dt|fieldset|figcaption|figure|"
    r"footer|form|frame|frameset|h[1-6]|head|header|hr|html|iframe|legend|li|"
    r"link|main|menu|menuitem|nav|noframes|ol|optgroup|option|p|param|section|"
    r"source|summary|table|tbody|td|tfoot|th|thead|title|tr|track|ul)"
    r"(?:[ \t]+|/?>|$)",
    flags=re.IGNORECASE,
)
_MARKDOWN_HTML_ATTRIBUTE_NAME = r"[A-Za-z_:][A-Za-z0-9_.:-]*"
_MARKDOWN_HTML_ATTRIBUTE_VALUE = (
    r'(?:[^\s"\'=<>`]+|\'[^\']*\'|"[^"]*")'
)
_MARKDOWN_HTML_ATTRIBUTE = (
    rf"(?:[ \t]+{_MARKDOWN_HTML_ATTRIBUTE_NAME}"
    rf"(?:[ \t]*=[ \t]*{_MARKDOWN_HTML_ATTRIBUTE_VALUE})?)"
)
_MARKDOWN_HTML_COMPLETE_OPEN_TAG_RE = re.compile(
    rf"^<(?P<tag>[A-Za-z][A-Za-z0-9-]*)"
    rf"{_MARKDOWN_HTML_ATTRIBUTE}*[ \t]*/?>[ \t]*$"
)
_MARKDOWN_HTML_COMPLETE_CLOSING_TAG_RE = re.compile(
    r"^</(?P<tag>[A-Za-z][A-Za-z0-9-]*)[ \t]*>[ \t]*$"
)
_MARKDOWN_ASCII_PUNCTUATION = frozenset(
    "!\"#$%&'()*+,-./:;<=>?@[\\]^_`{|}~"
)


@dataclass(frozen=True)
class _MarkdownListMarker:
    bullet_marker: str | None
    ordered_start: int | None
    ordered_delimiter: str | None
    has_content: bool
    indent_columns: int
    content_indent_columns: int

    @property
    def container_identity(self) -> tuple[str, str]:
        if self.bullet_marker is not None:
            return ("bullet", self.bullet_marker)
        assert self.ordered_delimiter is not None
        return ("ordered", self.ordered_delimiter)


@dataclass
class _MarkdownListFrame:
    container_identity: tuple[str, str]
    parent_content_indent_columns: int
    marker_indent_columns: int
    content_indent_columns: int
    paragraph_open: bool
    lazy_continuation_active: bool = False


@dataclass(frozen=True)
class _MarkdownListCandidate:
    parent_depth: int
    marker: _MarkdownListMarker
    marker_indent_columns: int
    content_indent_columns: int
    relative_line: str


@dataclass(frozen=True)
class ParsedLayoutItem:
    layout_id: str
    text: str
    char_start: int
    char_end: int
    page_number: int | None = None
    bbox: dict[str, Any] = field(default_factory=dict)
    coordinate_system: str = "text_flow_v1"
    region_type: str = "text"
    reading_order: int = 0
    confidence: float = 1.0
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass(frozen=True)
class ParsedStructureObject:
    structure_id: str
    object_type: str
    text: str
    char_start: int
    char_end: int
    title: str | None = None
    page_number: int | None = None
    bbox: dict[str, Any] = field(default_factory=dict)
    coordinate_system: str = "text_flow_v1"
    reading_order: int = 0
    parent_ref: str | None = None
    path: str | None = None
    confidence: float = 1.0
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass
class ParsedSection:
    title: str
    text: str
    page_number: int | None = None
    section: str | None = None
    metadata: dict[str, Any] = field(default_factory=dict)
    layout_items: list[ParsedLayoutItem] = field(default_factory=list)
    structure_objects: list[ParsedStructureObject] = field(default_factory=list)
    parser_metadata: dict[str, Any] = field(default_factory=dict)


def _unescaped_terminal_pipe(text: str) -> bool:
    if not text.endswith("|"):
        return False
    backslash_count = 0
    index = len(text) - 2
    while index >= 0 and text[index] == "\\":
        backslash_count += 1
        index -= 1
    return backslash_count % 2 == 0


def _markdown_leading_indent(line: str) -> tuple[int, int]:
    columns = 0
    character_count = 0
    for character in line or "":
        if character == " ":
            columns += 1
        elif character == "\t":
            columns += 4 - (columns % 4)
        else:
            break
        character_count += 1
    return columns, character_count


def _markdown_block_line_content(
    line: str,
    *,
    maximum_indent_columns: int = 3,
) -> str | None:
    indent_columns, indent_characters = _markdown_leading_indent(line)
    if indent_columns > maximum_indent_columns:
        return None
    return (line or "")[indent_characters:]


def _markdown_indented_code_line(line: str) -> bool:
    if not (line or "").strip():
        return False
    indent_columns, _indent_characters = _markdown_leading_indent(line)
    return indent_columns >= 4


def _split_markdown_pipe_row(line: str) -> tuple[str, ...] | None:
    content = _markdown_block_line_content(line)
    if content is None:
        return None
    stripped = content.strip(" \t")
    if not stripped or "\n" in stripped or "\r" in stripped:
        return None

    cells: list[str] = []
    current: list[str] = []
    separator_count = 0
    for index, character in enumerate(stripped):
        if character == "|":
            backslash_count = 0
            cursor = index - 1
            while cursor >= 0 and stripped[cursor] == "\\":
                backslash_count += 1
                cursor -= 1
            if backslash_count % 2 == 0:
                cells.append("".join(current).strip())
                current = []
                separator_count += 1
                continue
        current.append(character)
    cells.append("".join(current).strip())

    if separator_count == 0:
        return None
    if stripped.startswith("|"):
        cells = cells[1:]
    if _unescaped_terminal_pipe(stripped):
        cells = cells[:-1]
    return tuple(cells) if cells else None


def _markdown_table_body_row(line: str) -> tuple[str, ...] | None:
    """Return a GFM table body row, including a bare one-cell ragged row."""

    split_row = _split_markdown_pipe_row(line)
    if split_row is not None:
        return split_row
    content = _markdown_block_line_content(line)
    if content is None:
        return None
    stripped = content.strip(" \t")
    if not stripped or "\n" in stripped or "\r" in stripped:
        return None
    if _markdown_block_start_kind(line) is not None:
        return None
    return (stripped,)


def _markdown_table_alignment(delimiter_cell: str) -> str:
    left = delimiter_cell.startswith(":")
    right = delimiter_cell.endswith(":")
    if left and right:
        return "center"
    if left:
        return "left"
    if right:
        return "right"
    return "default"


def _markdown_pipe_table_descriptor(text: str) -> dict[str, Any] | None:
    lines = (text or "").splitlines()
    while lines and not lines[0].strip():
        lines.pop(0)
    while lines and not lines[-1].strip():
        lines.pop()
    if len(lines) < 2 or any(not line.strip() for line in lines):
        return None
    if _markdown_block_start_kind(lines[0]) is not None:
        return None
    header = _split_markdown_pipe_row(lines[0])
    delimiter = _split_markdown_pipe_row(lines[1])
    if (
        header is None
        or delimiter is None
        or len(header) != len(delimiter)
        or not any(cell for cell in header)
        or not all(_MARKDOWN_TABLE_DELIMITER_CELL_RE.fullmatch(cell) for cell in delimiter)
    ):
        return None
    if any(_markdown_block_start_kind(line) is not None for line in lines[2:]):
        return None
    source_body_rows = [_markdown_table_body_row(line) for line in lines[2:]]
    if any(row is None for row in source_body_rows):
        return None
    column_count = len(header)
    body_rows: list[tuple[str, ...]] = []
    source_body_column_counts: list[int] = []
    padded_missing_cell_counts: list[int] = []
    ignored_excess_cell_counts: list[int] = []
    for source_row in source_body_rows:
        assert source_row is not None
        source_column_count = len(source_row)
        source_body_column_counts.append(source_column_count)
        padded_missing_cell_counts.append(max(0, column_count - source_column_count))
        ignored_excess_cell_counts.append(max(0, source_column_count - column_count))
        body_rows.append(
            tuple(source_row[:column_count])
            + ("",) * max(0, column_count - source_column_count)
        )
    return {
        "column_count": column_count,
        "data_row_count": len(body_rows),
        "header_cells": tuple(header),
        "column_alignments": tuple(
            _markdown_table_alignment(cell) for cell in delimiter
        ),
        "body_rows": tuple(body_rows),
        "source_body_column_counts": tuple(source_body_column_counts),
        "padded_missing_cell_counts": tuple(padded_missing_cell_counts),
        "ignored_excess_cell_counts": tuple(ignored_excess_cell_counts),
        "row_normalization_protocol_version": (
            MARKDOWN_TABLE_ROW_NORMALIZATION_PROTOCOL_VERSION
        ),
        "target_spec_version": MARKDOWN_TARGET_SPEC_VERSION,
    }


def _markdown_pipe_table_shape(text: str) -> tuple[int, int] | None:
    descriptor = _markdown_pipe_table_descriptor(text)
    if descriptor is None:
        return None
    return int(descriptor["column_count"]), int(descriptor["data_row_count"])


def _markdown_fence_marker(line: str) -> tuple[str, int] | None:
    content = _markdown_block_line_content(line)
    if content is None:
        return None
    match = _MARKDOWN_FENCE_OPEN_RE.match(content)
    if match is None:
        return None
    marker = match.group(1)
    if marker.startswith("`") and "`" in match.group(2):
        return None
    return marker[0], len(marker)


def _markdown_fence_closes(line: str, marker: tuple[str, int]) -> bool:
    content = _markdown_block_line_content(line)
    if content is None:
        return False
    character, minimum_length = marker
    return bool(
        re.fullmatch(
            rf"{re.escape(character)}{{{minimum_length},}}[ \t]*",
            content,
        )
    )


def _markdown_atx_heading_title(line: str) -> str | None:
    if "\n" in (line or "") or "\r" in (line or ""):
        return None
    match = _MARKDOWN_ATX_OPEN_RE.fullmatch(line or "")
    if match is None:
        return None
    opening = match.group("opening")
    remainder = match.group("remainder")
    if len(opening) > 6 or (remainder and remainder[0] not in {" ", "\t"}):
        return None
    closing = _MARKDOWN_ATX_CLOSING_RE.search(remainder)
    if closing is not None:
        remainder = remainder[: closing.start()]
    return remainder.strip(" \t")


def _markdown_thematic_break(line_content: str) -> bool:
    compact = re.sub(r"[ \t]", "", line_content or "")
    return len(compact) >= 3 and compact[0] in {"*", "-", "_"} and len(
        set(compact)
    ) == 1


def _markdown_list_marker(line: str) -> _MarkdownListMarker | None:
    content = _markdown_block_line_content(line)
    if content is None:
        return None
    match = _MARKDOWN_LIST_MARKER_RE.fullmatch(content)
    if match is None:
        return None
    spacing = match.group("spacing") or ""
    item_content = match.group("content") or ""
    if item_content and not spacing:
        return None

    indent_columns, _indent_characters = _markdown_leading_indent(line)
    marker_width = 1
    ordered_start: int | None = None
    if match.group("number") is not None:
        ordered_start = int(match.group("number"))
        marker_width = len(match.group("number")) + 1
    content_indent_columns = indent_columns + marker_width
    for character in spacing:
        if character == "\t":
            content_indent_columns += 4 - (content_indent_columns % 4)
        else:
            content_indent_columns += 1
    return _MarkdownListMarker(
        bullet_marker=match.group("bullet"),
        ordered_start=ordered_start,
        ordered_delimiter=match.group("ordered_delimiter"),
        has_content=bool(item_content.strip(" \t")),
        indent_columns=indent_columns,
        content_indent_columns=content_indent_columns,
    )


def _markdown_list_interrupts_paragraph(line: str) -> bool:
    marker = _markdown_list_marker(line)
    if marker is None or not marker.has_content:
        return False
    return marker.ordered_start in {None, 1}


def _markdown_html_block_type(line_content: str) -> int | None:
    """Classify the seven HTML block types frozen by GFM 0.29."""

    content = line_content or ""
    if _MARKDOWN_HTML_TYPE_1_RE.match(content) is not None:
        return 1
    if content.startswith("<!--"):
        return 2
    if content.startswith("<?"):
        return 3
    if re.match(r"^<![A-Z]", content) is not None:
        return 4
    if content.startswith("<![CDATA["):
        return 5
    if _MARKDOWN_HTML_BLOCK_TAG_RE.match(content) is not None:
        return 6
    complete_open_tag = _MARKDOWN_HTML_COMPLETE_OPEN_TAG_RE.fullmatch(content)
    if complete_open_tag is not None:
        if str(complete_open_tag.group("tag") or "").lower() in {
            "pre",
            "script",
            "style",
        }:
            return None
        return 7
    if _MARKDOWN_HTML_COMPLETE_CLOSING_TAG_RE.fullmatch(content) is not None:
        return 7
    return None


def _markdown_html_block_start(line_content: str) -> bool:
    return _markdown_html_block_type(line_content) is not None


def _markdown_html_block_line_ends(line: str, html_block_type: int) -> bool:
    if html_block_type == 1:
        return _MARKDOWN_HTML_TYPE_1_END_RE.search(line or "") is not None
    if html_block_type == 2:
        return "-->" in (line or "")
    if html_block_type == 3:
        return "?>" in (line or "")
    if html_block_type == 4:
        return ">" in (line or "")
    if html_block_type == 5:
        return "]]>" in (line or "")
    return False


def _markdown_html_block_end(
    lines: list[tuple[int, int, str]],
    start_index: int,
    *,
    paragraph_open: bool = False,
) -> tuple[int, int] | None:
    """Return the exclusive end and GFM HTML type for a block opener."""

    if start_index >= len(lines):
        return None
    content = _markdown_block_line_content(lines[start_index][2])
    if content is None:
        return None
    html_block_type = _markdown_html_block_type(content)
    if html_block_type is None or (html_block_type == 7 and paragraph_open):
        return None
    cursor = start_index
    if html_block_type <= 5:
        while cursor < len(lines):
            cursor += 1
            if _markdown_html_block_line_ends(
                lines[cursor - 1][2],
                html_block_type,
            ):
                break
        return cursor, html_block_type
    while cursor < len(lines) and lines[cursor][2].strip():
        cursor += 1
    return cursor, html_block_type


def _markdown_link_reference_label_end(
    lines: list[tuple[int, int, str]],
    start_index: int,
) -> tuple[int, int, str] | None:
    """Return the closing-label line, post-colon offset, and raw label."""

    first_content = _markdown_block_line_content(lines[start_index][2])
    if first_content is None or not first_content.startswith("["):
        return None
    label_parts: list[str] = []
    label_character_count = 0
    line_index = start_index
    fragment = first_content
    cursor = 1
    while True:
        while cursor < len(fragment):
            character = fragment[cursor]
            if character == "\\" and cursor + 1 < len(fragment):
                label_parts.append(fragment[cursor : cursor + 2])
                label_character_count += 2
                cursor += 2
                continue
            if character == "[":
                return None
            if character == "]":
                if cursor + 1 >= len(fragment) or fragment[cursor + 1] != ":":
                    return None
                label = "".join(label_parts)
                if (
                    label_character_count > 999
                    or not any(not item.isspace() for item in label)
                ):
                    return None
                return line_index, cursor + 2, label
            label_parts.append(character)
            label_character_count += 1
            if label_character_count > 999:
                return None
            cursor += 1

        line_index += 1
        if line_index >= len(lines) or not lines[line_index][2].strip():
            return None
        continuation = _markdown_block_line_content(lines[line_index][2])
        if continuation is None:
            return None
        label_parts.append("\n")
        label_character_count += 1
        if label_character_count > 999:
            return None
        fragment = continuation
        cursor = 0


def _markdown_link_destination_end(text: str) -> int | None:
    """Return the destination end offset for the GFM link-reference grammar."""

    if not text:
        return None
    if text.startswith("<"):
        cursor = 1
        while cursor < len(text):
            character = text[cursor]
            if (
                character == "\\"
                and cursor + 1 < len(text)
                and text[cursor + 1] in _MARKDOWN_ASCII_PUNCTUATION
            ):
                cursor += 2
                continue
            if character == "<":
                return None
            if character == ">":
                return cursor + 1
            cursor += 1
        return None

    cursor = 0
    parenthesis_depth = 0
    while cursor < len(text) and text[cursor] not in {" ", "\t"}:
        character = text[cursor]
        if ord(character) < 32 or ord(character) == 127:
            return None
        if (
            character == "\\"
            and cursor + 1 < len(text)
            and text[cursor + 1] in _MARKDOWN_ASCII_PUNCTUATION
        ):
            cursor += 2
            continue
        if character == "(":
            parenthesis_depth += 1
            if parenthesis_depth > 32:
                return None
        elif character == ")":
            if parenthesis_depth == 0:
                return None
            parenthesis_depth -= 1
        cursor += 1
    if cursor == 0 or parenthesis_depth != 0:
        return None
    return cursor


def _markdown_link_reference_title_end(
    lines: list[tuple[int, int, str]],
    start_index: int,
    first_fragment: str,
) -> int | None:
    """Validate a possibly multiline GFM link-reference title."""

    if not first_fragment or first_fragment[0] not in {'"', "'", "("}:
        return None
    opener = first_fragment[0]
    closer = ")" if opener == "(" else opener
    fragment = first_fragment
    line_index = start_index
    fragment_cursor = 1
    while True:
        while fragment_cursor < len(fragment):
            character = fragment[fragment_cursor]
            if (
                character == "\\"
                and fragment_cursor + 1 < len(fragment)
                and fragment[fragment_cursor + 1] in _MARKDOWN_ASCII_PUNCTUATION
            ):
                fragment_cursor += 2
                continue
            if opener == "(" and character == "(":
                return None
            if character == closer:
                if fragment[fragment_cursor + 1 :].strip(" \t"):
                    return None
                return line_index + 1
            fragment_cursor += 1

        line_index += 1
        if line_index >= len(lines) or not lines[line_index][2].strip():
            return None
        continuation = _markdown_block_line_content(lines[line_index][2])
        if continuation is None:
            return None
        fragment = continuation
        fragment_cursor = 0


def _markdown_link_reference_definition_end(
    lines: list[tuple[int, int, str]],
    start_index: int,
) -> int | None:
    """Return the exclusive end of a valid GFM link reference definition.

    The destination may start on the following non-blank line and a title may
    span non-blank lines. Speculative continuation never has block authority:
    the caller receives a boundary only after the complete grammar validates.
    """

    if start_index >= len(lines):
        return None
    label_end = _markdown_link_reference_label_end(lines, start_index)
    if label_end is None:
        return None
    label_line_index, post_colon_offset, _label = label_end
    label_line_content = _markdown_block_line_content(lines[label_line_index][2])
    assert label_line_content is not None

    destination_line_index = label_line_index
    destination_fragment = label_line_content[post_colon_offset:].lstrip(" \t")
    if not destination_fragment:
        destination_line_index += 1
        if (
            destination_line_index >= len(lines)
            or not lines[destination_line_index][2].strip()
        ):
            return None
        destination_content = _markdown_block_line_content(
            lines[destination_line_index][2]
        )
        if destination_content is None:
            return None
        destination_fragment = destination_content.lstrip(" \t")

    destination_end = _markdown_link_destination_end(destination_fragment)
    if destination_end is None:
        return None
    destination_tail = destination_fragment[destination_end:]
    if destination_tail:
        if destination_tail[0] not in {" ", "\t"}:
            return None
        title_fragment = destination_tail.lstrip(" \t")
        if title_fragment:
            return _markdown_link_reference_title_end(
                lines,
                destination_line_index,
                title_fragment,
            )

    definition_end = destination_line_index + 1
    if definition_end >= len(lines) or not lines[definition_end][2].strip():
        return definition_end
    possible_title = _markdown_block_line_content(lines[definition_end][2])
    if possible_title is None:
        return definition_end
    possible_title = possible_title.lstrip(" \t")
    if not possible_title or possible_title[0] not in {'"', "'", "("}:
        return definition_end
    optional_title_end = _markdown_link_reference_title_end(
        lines,
        definition_end,
        possible_title,
    )
    return optional_title_end or definition_end


def _markdown_link_reference_definition_start(line_content: str) -> bool:
    content = line_content or ""
    single_line = [(0, len(content), content)]
    return _markdown_link_reference_definition_end(single_line, 0) == 1


def _markdown_block_start_kind(line: str) -> str | None:
    """Classify non-table GFM block starts before parsing inline pipe cells."""

    if not (line or "").strip():
        return "blank"
    if _markdown_indented_code_line(line):
        return "indented_code"
    content = _markdown_block_line_content(line)
    if content is None:
        return None
    if _markdown_fence_marker(line) is not None:
        return "fenced_code"
    if content.startswith(">"):
        return "block_quote"
    if _markdown_atx_heading_title(line) is not None:
        return "atx_heading"
    if _markdown_thematic_break(content):
        return "thematic_break"
    if _markdown_list_marker(line) is not None:
        return "list"
    if _markdown_html_block_start(content):
        return "html_block"
    if _markdown_link_reference_definition_start(content):
        return "link_reference_definition"
    return None


def _markdown_strip_container_indent(
    line: str,
    indent_columns: int,
) -> str | None:
    """Consume a container prefix and preserve any unconsumed tab columns."""

    if indent_columns < 0:
        return None
    consumed_columns = 0
    cursor = 0
    while consumed_columns < indent_columns and cursor < len(line or ""):
        character = line[cursor]
        if character == " ":
            width = 1
        elif character == "\t":
            width = 4 - (consumed_columns % 4)
        else:
            return None
        cursor += 1
        consumed_columns += width
    if consumed_columns < indent_columns:
        return None
    residual_indent_columns = consumed_columns - indent_columns
    while cursor < len(line or ""):
        character = line[cursor]
        if character == " ":
            width = 1
        elif character == "\t":
            width = 4 - (consumed_columns % 4)
        else:
            break
        cursor += 1
        consumed_columns += width
        residual_indent_columns += width
    return " " * residual_indent_columns + (line or "")[cursor:]


def _markdown_list_content_block_start_kind(
    line: str,
    *,
    content_indent_columns: int,
) -> str | None:
    """Classify a leaf only after consuming the active list-item prefix."""

    block_start = _markdown_list_content_block_start(
        line,
        content_indent_columns=content_indent_columns,
    )
    if block_start is None:
        return None
    return block_start[0]


def _markdown_list_content_block_start(
    line: str,
    *,
    content_indent_columns: int,
) -> tuple[str | None, str] | None:
    """Return one container-relative block classification and its residual."""

    content = _markdown_strip_container_indent(line, content_indent_columns)
    if content is None:
        return None
    return _markdown_block_start_kind(content), content


def _markdown_list_candidate(
    line: str,
    frames: list[_MarkdownListFrame],
) -> _MarkdownListCandidate | None:
    """Return the deepest list marker whose parent prefix is active."""

    parent_prefixes = [(-1, 0)] + [
        (depth, frame.content_indent_columns)
        for depth, frame in enumerate(frames)
    ]
    for parent_depth, parent_prefix in reversed(parent_prefixes):
        relative_line = _markdown_strip_container_indent(line, parent_prefix)
        if relative_line is None:
            continue
        if _markdown_block_start_kind(relative_line) != "list":
            continue
        marker = _markdown_list_marker(relative_line)
        if marker is None:
            continue
        return _MarkdownListCandidate(
            parent_depth=parent_depth,
            marker=marker,
            marker_indent_columns=parent_prefix + marker.indent_columns,
            content_indent_columns=(
                parent_prefix + marker.content_indent_columns
            ),
            relative_line=relative_line,
        )
    return None


def _markdown_list_lazy_continuation(
    *,
    frame: _MarkdownListFrame,
    relative_kind: str | None,
    relative_line: str,
) -> bool:
    if not frame.paragraph_open:
        return False
    if relative_kind in {None, "indented_code"}:
        return True
    if relative_kind != "html_block":
        return False
    html_content = _markdown_block_line_content(relative_line)
    return (
        frame.lazy_continuation_active
        and html_content is not None
        and _markdown_html_block_type(html_content) == 7
    )


def _flow_line_spans(text: str) -> list[tuple[int, int, str]]:
    lines: list[tuple[int, int, str]] = []
    cursor = 0
    for raw_line in (text or "").splitlines(keepends=True):
        line = raw_line.rstrip("\r\n")
        lines.append((cursor, cursor + len(line), line))
        cursor += len(raw_line)
    if text and not lines:
        lines.append((0, len(text), text))
    return lines


def _markdown_table_end(
    lines: list[tuple[int, int, str]],
    start_index: int,
) -> int | None:
    if start_index + 1 >= len(lines):
        return None
    if _markdown_block_start_kind(lines[start_index][2]) is not None:
        return None
    header = _split_markdown_pipe_row(lines[start_index][2])
    delimiter = _split_markdown_pipe_row(lines[start_index + 1][2])
    if (
        header is None
        or delimiter is None
        or len(header) != len(delimiter)
        or not any(cell for cell in header)
        or not all(_MARKDOWN_TABLE_DELIMITER_CELL_RE.fullmatch(cell) for cell in delimiter)
    ):
        return None

    cursor = start_index + 2
    while cursor < len(lines):
        if _markdown_link_reference_definition_end(lines, cursor) is not None:
            break
        if _markdown_block_start_kind(lines[cursor][2]) is not None:
            break
        row = _markdown_table_body_row(lines[cursor][2])
        if row is None:
            break
        cursor += 1
    return cursor


def _markdown_indented_code_end(
    lines: list[tuple[int, int, str]],
    start_index: int,
) -> int | None:
    if (
        start_index >= len(lines)
        or not _markdown_indented_code_line(lines[start_index][2])
    ):
        return None
    cursor = start_index
    last_indented_line_end = start_index
    while cursor < len(lines):
        line = lines[cursor][2]
        if not line.strip():
            cursor += 1
            continue
        if not _markdown_indented_code_line(line):
            break
        last_indented_line_end = cursor + 1
        cursor += 1
    return last_indented_line_end


def _markdown_list_block_end(
    lines: list[tuple[int, int, str]],
    start_index: int,
) -> int | None:
    """Return one top-level list span using an explicit container stack."""

    marker = _markdown_list_marker(lines[start_index][2])
    if marker is None:
        return None
    frames = [
        _MarkdownListFrame(
            container_identity=marker.container_identity,
            parent_content_indent_columns=0,
            marker_indent_columns=marker.indent_columns,
            content_indent_columns=marker.content_indent_columns,
            paragraph_open=marker.has_content,
        )
    ]
    cursor = start_index + 1
    while cursor < len(lines):
        line = lines[cursor][2]
        if not line.strip():
            blank_start = cursor
            while cursor < len(lines) and not lines[cursor][2].strip():
                cursor += 1
            if cursor >= len(lines):
                return blank_start
            for frame in frames:
                frame.paragraph_open = False
                frame.lazy_continuation_active = False
            lookahead_line = lines[cursor][2]
            lookahead_candidate = _markdown_list_candidate(
                lookahead_line,
                frames,
            )
            if lookahead_candidate is not None:
                target_depth = lookahead_candidate.parent_depth + 1
                if not (
                    target_depth == 0
                    and lookahead_candidate.marker.container_identity
                    != frames[0].container_identity
                ):
                    continue
            if any(
                _markdown_strip_container_indent(
                    lookahead_line,
                    frame.content_indent_columns,
                )
                is not None
                for frame in frames
            ):
                continue
            return blank_start

        candidate = _markdown_list_candidate(line, frames)
        if candidate is not None:
            target_depth = candidate.parent_depth + 1
            starts_new_container = (
                target_depth == len(frames)
                or (
                    target_depth > 0
                    and candidate.marker.container_identity
                    != frames[target_depth].container_identity
                )
            )
            if (
                starts_new_container
                and frames[-1].paragraph_open
                and not _markdown_list_interrupts_paragraph(
                    candidate.relative_line
                )
            ):
                frames[-1].paragraph_open = True
                frames[-1].lazy_continuation_active = (
                    target_depth < len(frames)
                )
                cursor += 1
                continue
            elif (
                target_depth == 0
                and candidate.marker.container_identity
                != frames[0].container_identity
            ):
                break
            else:
                if target_depth > 0:
                    parent = frames[target_depth - 1]
                    parent.paragraph_open = False
                    parent.lazy_continuation_active = False
                new_frame = _MarkdownListFrame(
                    container_identity=candidate.marker.container_identity,
                    parent_content_indent_columns=(
                        0
                        if candidate.parent_depth < 0
                        else frames[candidate.parent_depth].content_indent_columns
                    ),
                    marker_indent_columns=candidate.marker_indent_columns,
                    content_indent_columns=candidate.content_indent_columns,
                    paragraph_open=candidate.marker.has_content,
                )
                if target_depth < len(frames):
                    frames[target_depth:] = [new_frame]
                else:
                    frames.append(new_frame)
                cursor += 1
                continue

        matched_depth = -1
        relative_line = line
        for depth in range(len(frames) - 1, -1, -1):
            possible_relative_line = _markdown_strip_container_indent(
                line,
                frames[depth].content_indent_columns,
            )
            if possible_relative_line is not None:
                matched_depth = depth
                relative_line = possible_relative_line
                break
        relative_kind = _markdown_block_start_kind(relative_line)
        innermost = frames[-1]
        if matched_depth < len(frames) - 1 and _markdown_list_lazy_continuation(
            frame=innermost,
            relative_kind=relative_kind,
            relative_line=relative_line,
        ):
            innermost.paragraph_open = True
            innermost.lazy_continuation_active = True
            cursor += 1
            continue
        if matched_depth < 0:
            break
        if matched_depth < len(frames) - 1:
            frames[matched_depth + 1 :] = []
        frame = frames[-1]
        if relative_kind is None:
            frame.paragraph_open = True
        elif relative_kind == "indented_code" and frame.paragraph_open:
            # Indented code cannot interrupt the current item paragraph.
            frame.paragraph_open = True
        elif relative_kind == "html_block":
            html_content = _markdown_block_line_content(relative_line)
            frame.paragraph_open = bool(
                frame.paragraph_open
                and html_content is not None
                and _markdown_html_block_type(html_content) == 7
            )
        elif relative_kind == "list" and frame.paragraph_open:
            frame.paragraph_open = not _markdown_list_interrupts_paragraph(
                relative_line
            )
        else:
            frame.paragraph_open = False
        frame.lazy_continuation_active = False
        cursor += 1
    return cursor


def _plain_flow_spans(text: str) -> list[tuple[int, int, str | None]]:
    lines = _flow_line_spans(text)
    spans: list[tuple[int, int, str | None]] = []
    cursor = 0
    while cursor < len(lines):
        if not lines[cursor][2].strip():
            cursor += 1
            continue
        start = cursor
        cursor += 1
        while cursor < len(lines) and lines[cursor][2].strip():
            cursor += 1
        spans.append((lines[start][0], lines[cursor - 1][1], None))
    return spans


def _markdown_flow_spans(text: str) -> list[tuple[int, int, str | None]]:
    lines = _flow_line_spans(text)
    spans: list[tuple[int, int, str | None]] = []
    cursor = 0
    while cursor < len(lines):
        if not lines[cursor][2].strip():
            cursor += 1
            continue

        block_start_kind = _markdown_block_start_kind(lines[cursor][2])
        if block_start_kind == "indented_code":
            indented_code_end = _markdown_indented_code_end(lines, cursor)
            assert indented_code_end is not None
            spans.append(
                (
                    lines[cursor][0],
                    lines[indented_code_end - 1][1],
                    "code_block",
                )
            )
            cursor = indented_code_end
            continue

        if block_start_kind == "fenced_code":
            fence = _markdown_fence_marker(lines[cursor][2])
            assert fence is not None
            start = cursor
            cursor += 1
            while cursor < len(lines):
                if _markdown_fence_closes(lines[cursor][2], fence):
                    cursor += 1
                    break
                cursor += 1
            spans.append((lines[start][0], lines[cursor - 1][1], "code_block"))
            continue

        html_block = _markdown_html_block_end(lines, cursor)
        if html_block is not None:
            html_block_end, _html_block_type = html_block
            spans.append(
                (
                    lines[cursor][0],
                    lines[html_block_end - 1][1],
                    None,
                )
            )
            cursor = html_block_end
            continue

        table_end = _markdown_table_end(lines, cursor)
        if table_end is not None:
            spans.append((lines[cursor][0], lines[table_end - 1][1], "table"))
            cursor = table_end
            continue

        if block_start_kind == "list":
            list_end = _markdown_list_block_end(lines, cursor)
            assert list_end is not None
            spans.append((lines[cursor][0], lines[list_end - 1][1], "list"))
            cursor = list_end
            continue

        if block_start_kind in {"atx_heading", "thematic_break"}:
            spans.append((lines[cursor][0], lines[cursor][1], None))
            cursor += 1
            continue

        start = cursor
        cursor += 1
        while cursor < len(lines) and lines[cursor][2].strip():
            if _markdown_table_end(lines, cursor) is not None:
                break
            next_block_start_kind = _markdown_block_start_kind(lines[cursor][2])
            if next_block_start_kind == "indented_code":
                cursor += 1
                continue
            if next_block_start_kind == "link_reference_definition":
                cursor += 1
                continue
            if next_block_start_kind == "list":
                if not _markdown_list_interrupts_paragraph(lines[cursor][2]):
                    cursor += 1
                    continue
                break
            if next_block_start_kind == "html_block":
                html_content = _markdown_block_line_content(lines[cursor][2])
                if (
                    html_content is not None
                    and _markdown_html_block_type(html_content) == 7
                ):
                    cursor += 1
                    continue
                break
            if next_block_start_kind is not None:
                if (
                    block_start_kind == next_block_start_kind
                    and block_start_kind in {"block_quote", "list"}
                ):
                    cursor += 1
                    continue
                break
            cursor += 1
        spans.append((lines[start][0], lines[cursor - 1][1], None))
    return spans


def _block_object_type(text: str, *, content_kind: str = "text") -> str:
    stripped = (text or "").strip()
    lowered = stripped.lower()
    if (
        content_kind == "code"
        or lowered.startswith("[code cell]")
        or _markdown_fence_marker(stripped.splitlines()[0] if stripped else "") is not None
    ):
        return "code_block"
    if _markdown_pipe_table_shape(text) is not None or (
        content_kind != "markdown" and "\t" in stripped
    ):
        return "table"
    if _detect_formula(stripped) or re.search(r"\$\$.*?\$\$", stripped, flags=re.DOTALL):
        return "formula"
    if re.match(r"(?i)^(?:figure|fig\.|table|caption|图|表)\s*[:：\d.-]", stripped):
        return "caption"
    if re.match(r"^\s*(?:[-*+]\s+|\d+[.)]\s+)", stripped):
        return "list"
    return "paragraph"


def _flow_artifacts(
    text: str,
    *,
    artifact_prefix: str,
    page_number: int | None,
    content_kind: str,
    coordinate_system: str = "text_flow_v1",
    bbox: dict[str, Any] | None = None,
    source: str,
    path_prefix: str | None = None,
) -> tuple[list[ParsedLayoutItem], list[ParsedStructureObject]]:
    spans = (
        _markdown_flow_spans(text)
        if content_kind == "markdown"
        else _plain_flow_spans(text)
    )
    layouts: list[ParsedLayoutItem] = []
    objects: list[ParsedStructureObject] = []
    for order, (start, end, detected_type) in enumerate(spans):
        block_text = text[start:end]
        object_type = detected_type or _block_object_type(
            block_text,
            content_kind=content_kind,
        )
        item_bbox = dict(bbox or {})
        common_metadata = {
            "parser_source": source,
            "native_geometry": bool(item_bbox),
            "layout_protocol_version": PARSER_LAYOUT_PROTOCOL_VERSION,
            "flow_block_protocol_version": FLOW_BLOCK_PROTOCOL_VERSION,
        }
        if content_kind == "markdown":
            common_metadata["block_start_protocol_version"] = (
                MARKDOWN_BLOCK_START_PROTOCOL_VERSION
            )
            common_metadata.update(
                {
                    "html_block_protocol_version": (
                        MARKDOWN_HTML_BLOCK_PROTOCOL_VERSION
                    ),
                    "link_reference_protocol_version": (
                        MARKDOWN_LINK_REFERENCE_PROTOCOL_VERSION
                    ),
                    "target_spec_version": MARKDOWN_TARGET_SPEC_VERSION,
                }
            )
        table_descriptor = (
            _markdown_pipe_table_descriptor(block_text)
            if object_type == "table" and content_kind == "markdown"
            else None
        )
        if table_descriptor is not None:
            common_metadata.update(
                {
                    "structure_protocol_version": MARKDOWN_PIPE_TABLE_PROTOCOL_VERSION,
                    "column_count": table_descriptor["column_count"],
                    "data_row_count": table_descriptor["data_row_count"],
                    "column_alignments": list(
                        table_descriptor["column_alignments"]
                    ),
                    "source_body_column_counts": list(
                        table_descriptor["source_body_column_counts"]
                    ),
                    "padded_missing_cell_counts": list(
                        table_descriptor["padded_missing_cell_counts"]
                    ),
                    "ignored_excess_cell_counts": list(
                        table_descriptor["ignored_excess_cell_counts"]
                    ),
                    "row_normalization_protocol_version": (
                        table_descriptor["row_normalization_protocol_version"]
                    ),
                    "target_spec_version": table_descriptor["target_spec_version"],
                }
            )
        if content_kind == "markdown":
            first_line = block_text.splitlines()[0] if block_text else ""
            first_line_content = _markdown_block_line_content(first_line)
            html_block_type = (
                _markdown_html_block_type(first_line_content)
                if first_line_content is not None
                else None
            )
            if html_block_type is not None:
                common_metadata["html_block_type"] = html_block_type
        if (
            object_type == "code_block"
            and content_kind == "markdown"
            and _markdown_indented_code_line(block_text.splitlines()[0])
        ):
            source_indent_columns = [
                _markdown_leading_indent(line)[0]
                for line in block_text.splitlines()
                if line.strip()
            ]
            common_metadata.update(
                {
                    "structure_protocol_version": (
                        MARKDOWN_INDENTED_CODE_PROTOCOL_VERSION
                    ),
                    "indentation_protocol": "tab_stop_4_remove_4_columns_v1",
                    "source_indent_columns": source_indent_columns,
                    "source_span_protocol": "raw_markdown_char_span_v1",
                }
            )
        structure_path = (
            f"{path_prefix} / {object_type}:{order + 1}"
            if path_prefix
            else f"{object_type}:{order + 1}"
        )
        layouts.append(
            ParsedLayoutItem(
                layout_id=f"{artifact_prefix}:layout:{order}",
                text=block_text,
                char_start=start,
                char_end=end,
                page_number=page_number,
                bbox=item_bbox,
                coordinate_system=coordinate_system,
                region_type=object_type,
                reading_order=order,
                metadata=common_metadata,
            )
        )
        objects.append(
            ParsedStructureObject(
                structure_id=f"{artifact_prefix}:structure:{order}",
                object_type=object_type,
                text=block_text,
                char_start=start,
                char_end=end,
                title=re.sub(r"\s+", " ", block_text).strip()[:120] or object_type,
                page_number=page_number,
                bbox=item_bbox,
                coordinate_system=coordinate_system,
                reading_order=order,
                path=structure_path,
                metadata=common_metadata,
            )
        )
    return layouts, objects


def _parser_path(source: FrozenSourceSnapshot) -> Path:
    return source.canonical_path


def _test_only_frozen_source(path: Path) -> FrozenSourceSnapshot:
    """Freeze parser-only test input; unavailable without the explicit adapter."""

    if _TEST_DURABILITY_ADAPTER.get() is None:
        raise TypeError("parse_document requires FrozenSourceSnapshot")
    lexical = Path(os.path.abspath(path))
    before = lexical.lstat()
    if not stat.S_ISREG(before.st_mode) or int(before.st_nlink) != 1:
        raise TypeError("Parser test source must be one regular file")
    limit = int(get_settings().upload_max_bytes)
    if int(before.st_size) > limit:
        raise ValueError("Parser test source exceeds upload hard limit")
    with lexical.open("rb") as handle:
        descriptor = os.fstat(handle.fileno())
        content = handle.read(limit + 1)
        after_descriptor = os.fstat(handle.fileno())
    after = lexical.lstat()
    if (
        len(content) > limit
        or len(content) != int(before.st_size)
        or (int(before.st_dev), int(before.st_ino), int(before.st_size))
        != (int(descriptor.st_dev), int(descriptor.st_ino), int(descriptor.st_size))
        or (int(descriptor.st_dev), int(descriptor.st_ino), int(descriptor.st_size))
        != (
            int(after_descriptor.st_dev),
            int(after_descriptor.st_ino),
            int(after_descriptor.st_size),
        )
        or (
            int(before.st_dev),
            int(before.st_ino),
            int(before.st_size),
            int(before.st_mtime_ns),
            int(before.st_ctime_ns),
            int(before.st_nlink),
        )
        != (
            int(after.st_dev),
            int(after.st_ino),
            int(after.st_size),
            int(after.st_mtime_ns),
            int(after.st_ctime_ns),
            int(after.st_nlink),
        )
    ):
        raise RuntimeError("Parser test source identity changed during freeze")
    validated = _validate_seekable_upload_content(
        io.BytesIO(content),
        filename=lexical.name,
        max_bytes=limit,
        expected_checksum=hashlib.sha256(content).hexdigest(),
    )
    return FrozenSourceSnapshot(
        canonical_path=lexical,
        checksum=validated.checksum,
        size_bytes=validated.size_bytes,
        content_kind=validated.content_kind,
        suffix=validated.suffix,
        identity=VerifiedSourceIdentity(
            protocol_version="explicit_parser_test_freeze_v1",
            root_device_id=int(before.st_dev),
            root_inode=int(before.st_ino),
            device_id=int(before.st_dev),
            inode=int(before.st_ino),
            size_bytes=int(before.st_size),
            mtime_ns=int(before.st_mtime_ns),
            ctime_ns=int(before.st_ctime_ns),
            link_count=int(before.st_nlink),
        ),
        content_bytes=content,
        protocol_version=FROZEN_SOURCE_SNAPSHOT_PROTOCOL_VERSION,
    )


def _require_frozen_source(source: FrozenSourceSnapshot) -> FrozenSourceSnapshot:
    if isinstance(source, FrozenSourceSnapshot):
        return validate_frozen_source_snapshot(source)
    # Existing parser unit tests run only under the explicit durability
    # adapter. Production callers cannot activate this branch.
    if isinstance(source, Path):
        return _test_only_frozen_source(source)
    raise TypeError("parse_document requires FrozenSourceSnapshot")


def detect_source_type(source: FrozenSourceSnapshot) -> str:
    return source_type_from_path(_parser_path(source))


CONTROL_CHAR_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]")
MOJIBAKE_MARKERS = (
    "\ufffd",
    "\u00c3",
    "\u00c2",
    "\u00e2",
    "\u9208",
    "\u9365",
    "\u9429",
    "\u95b3",
    "\u951f",
    "\u7d34",
    "\u6d93",
    "\u934f",
    "\u704f",
    "\u93c4",
    "\u9a9e",
)
MOJIBAKE_MARKER_RE = re.compile("|".join(re.escape(marker) for marker in MOJIBAKE_MARKERS))
LATIN_WORD_HYPHEN_BREAK_RE = re.compile(r"(?<=[A-Za-z])-\n(?=[a-z])")
SOFT_SINGLE_NEWLINE_RE = re.compile(r"(?<![\n.!?:;。！？：；])\n(?!\s*(?:\n|[#>*+\-]|\d+[.)]))")
TEXT_ENCODINGS = ("utf-8-sig", "utf-8", "gb18030", "big5", "cp1252", "latin-1")


def _safe_import_ftfy():
    try:
        import ftfy

        return ftfy
    except Exception:
        return None


def _detect_encoding(raw: bytes) -> tuple[str | None, float | None]:
    try:
        from charset_normalizer import from_bytes

        match = from_bytes(raw).best()
        if match is not None and match.encoding:
            coherence = getattr(match, "percent_coherence", None)
            return str(match.encoding), float(coherence) if coherence is not None else None
    except Exception:
        pass
    return None, None


def _mojibake_score(text: str) -> float:
    if not text:
        return 0.0
    markers = len(MOJIBAKE_MARKER_RE.findall(text))
    replacement = text.count("\ufffd")
    controls = len(CONTROL_CHAR_RE.findall(text))
    lossy_latin_questions = len(re.findall(r"(?<=[A-Za-z])\?{2,}(?=\b|[\s.,;:])", text))
    return (markers * 2.0 + replacement * 3.0 + controls * 2.0 + lossy_latin_questions * 2.5) / max(len(text), 1)


def _repair_lossy_latin_question_marks(text: str) -> str:
    # Some extractors replace non-ASCII Latin letters with repeated question
    # marks. The original bytes are gone, so keep this limited to common words.
    return re.sub(r"\b([Cc])af\?{2,}(?=\s|$|[.,;:])", lambda match: f"{match.group(1)}af\u00e9", text)


def _repair_mojibake_candidate(text: str) -> tuple[str, bool]:
    original_score = _mojibake_score(text)
    if original_score <= 0:
        return text, False

    candidates = [text]
    ftfy = _safe_import_ftfy()
    if ftfy is not None:
        try:
            candidates.append(ftfy.fix_text(text, normalization="NFC"))
        except Exception:
            pass
    candidates.append(_repair_lossy_latin_question_marks(text))

    # Typical PDF / web extraction failure: UTF-8 bytes were decoded as a
    # legacy code page, producing CJK-looking garbage such as "閺嶇绺?.
    for encoding in ("gb18030", "big5", "cp1252", "latin-1"):
        try:
            candidates.append(text.encode(encoding, errors="strict").decode("utf-8", errors="strict"))
        except Exception:
            continue

    best = min(candidates, key=_mojibake_score)
    best_score = _mojibake_score(best)
    if best != text and best_score + 0.002 < original_score:
        return best, True
    return text, False


def decode_text_bytes(raw: bytes) -> tuple[str, dict[str, Any]]:
    metadata: dict[str, Any] = {}
    detected, coherence = _detect_encoding(raw)
    if detected:
        metadata["encoding_detected"] = detected
    if coherence is not None:
        metadata["encoding_coherence"] = round(coherence, 3)

    preferred = []
    detected_usable = detected and (coherence is None or coherence >= 20.0 or detected.lower().replace("_", "-") in TEXT_ENCODINGS)
    if detected_usable:
        preferred.append(detected)
    preferred.extend(encoding for encoding in TEXT_ENCODINGS if encoding.lower() not in {item.lower() for item in preferred})

    last_error: Exception | None = None
    for encoding in preferred:
        try:
            text = raw.decode(encoding)
            metadata.setdefault("encoding_used", encoding)
            return text, metadata
        except (LookupError, UnicodeDecodeError) as exc:
            last_error = exc
            continue
    metadata["encoding_used"] = "utf-8-ignore"
    if last_error is not None:
        metadata["encoding_error"] = str(last_error)
    return raw.decode("utf-8", errors="ignore"), metadata


def clean_extracted_text(
    text: str,
    *,
    source_type: str | None = None,
    preserve_horizontal_whitespace: bool = False,
    preserve_boundary_whitespace: bool = False,
) -> tuple[str, dict[str, Any]]:
    flags: list[str] = []
    original = text or ""
    cleaned = original.replace("\r\n", "\n").replace("\r", "\n").replace("\ufeff", "")
    if cleaned != original:
        flags.append("normalized_line_endings_or_bom")
    before_controls = cleaned
    cleaned = CONTROL_CHAR_RE.sub("", cleaned)
    if cleaned != before_controls:
        flags.append("removed_control_chars")
    cleaned = unicodedata.normalize("NFC", cleaned)

    repaired, repaired_mojibake = _repair_mojibake_candidate(cleaned)
    if repaired_mojibake:
        cleaned = repaired
        flags.append("mojibake_repaired")

    if source_type in {"pdf", "image", "ocr"}:
        before_layout = cleaned
        cleaned = LATIN_WORD_HYPHEN_BREAK_RE.sub("", cleaned)
        cleaned = SOFT_SINGLE_NEWLINE_RE.sub(" ", cleaned)
        if cleaned != before_layout:
            flags.append("normalized_pdf_ocr_linebreaks")

    before_space = cleaned
    if not preserve_horizontal_whitespace:
        cleaned = re.sub(r"[ \t]{2,}", " ", cleaned)
    cleaned = re.sub(r"\n{4,}", "\n\n\n", cleaned)
    if not preserve_boundary_whitespace:
        cleaned = cleaned.strip()
    comparison_before = (
        before_space
        if preserve_boundary_whitespace
        else before_space.strip()
    )
    if cleaned != comparison_before:
        flags.append("normalized_whitespace")

    metadata = {
        "text_cleaning_flags": sorted(set(flags)),
        "mojibake_repaired": repaired_mojibake,
        "mojibake_score_before": round(_mojibake_score(original), 6),
        "mojibake_score_after": round(_mojibake_score(cleaned), 6),
    }
    return cleaned, metadata


def load_text_with_metadata(
    source: FrozenSourceSnapshot,
    *,
    preserve_horizontal_whitespace: bool = False,
    preserve_boundary_whitespace: bool = False,
) -> tuple[str, dict[str, Any]]:
    source = _require_frozen_source(source)
    decoded, metadata = decode_text_bytes(source.content_bytes)
    cleaned, cleaning = clean_extracted_text(
        decoded,
        source_type=detect_source_type(source),
        preserve_horizontal_whitespace=preserve_horizontal_whitespace,
        preserve_boundary_whitespace=preserve_boundary_whitespace,
    )
    return cleaned, {**metadata, **cleaning}


def load_text(path: Path) -> str:
    return load_text_with_metadata(path)[0]


def _remap_parsed_artifacts(
    artifacts: list[ParsedLayoutItem] | list[ParsedStructureObject],
    *,
    original_text: str,
    cleaned_text: str,
    source_type: str,
) -> list[ParsedLayoutItem] | list[ParsedStructureObject]:
    remapped: list[ParsedLayoutItem | ParsedStructureObject] = []
    cursor = 0
    original_length = max(1, len(original_text))
    for artifact in sorted(artifacts, key=lambda item: (item.reading_order, item.char_start, item.char_end)):
        fragment, _cleaning = clean_extracted_text(
            artifact.text,
            source_type=source_type,
            preserve_horizontal_whitespace=source_type == "markdown",
            preserve_boundary_whitespace=source_type == "markdown",
        )
        start = cleaned_text.find(fragment, cursor) if fragment else -1
        remap_method = "cleaned_text_match"
        confidence = float(artifact.confidence)
        if start < 0 and fragment:
            start = cleaned_text.find(fragment)
        if start < 0:
            start = min(len(cleaned_text), round(max(0, artifact.char_start) / original_length * len(cleaned_text)))
            end = min(len(cleaned_text), round(max(artifact.char_start, artifact.char_end) / original_length * len(cleaned_text)))
            fragment = cleaned_text[start:end]
            remap_method = "proportional_fallback"
            confidence *= 0.7
        else:
            end = min(len(cleaned_text), start + len(fragment))
        if end <= start:
            continue
        metadata = {
            **dict(artifact.metadata or {}),
            "span_remap_method": remap_method,
            "original_char_span": [artifact.char_start, artifact.char_end],
        }
        remapped.append(
            replace(
                artifact,
                text=fragment,
                char_start=start,
                char_end=end,
                confidence=round(confidence, 6),
                metadata=metadata,
            )
        )
        cursor = max(cursor, end)
    return remapped


def _clean_section(section: ParsedSection, source_type: str, common_metadata: dict[str, Any] | None = None) -> ParsedSection:
    text, cleaning = clean_extracted_text(
        section.text,
        source_type=source_type,
        preserve_horizontal_whitespace=source_type == "markdown",
        preserve_boundary_whitespace=source_type == "markdown",
    )
    title, title_cleaning = clean_extracted_text(section.title, source_type=source_type)
    section_label = section.section
    if section_label is not None:
        section_label, _ = clean_extracted_text(section_label, source_type=source_type)
    flags = sorted(set((cleaning.get("text_cleaning_flags") or []) + (title_cleaning.get("text_cleaning_flags") or [])))
    metadata = {
        **(common_metadata or {}),
        **section.metadata,
        **cleaning,
        "text_cleaning_flags": flags,
        "mojibake_repaired": bool(cleaning.get("mojibake_repaired") or title_cleaning.get("mojibake_repaired")),
    }
    layout_items = _remap_parsed_artifacts(
        section.layout_items,
        original_text=section.text,
        cleaned_text=text,
        source_type=source_type,
    )
    structure_objects = _remap_parsed_artifacts(
        section.structure_objects,
        original_text=section.text,
        cleaned_text=text,
        source_type=source_type,
    )
    parser_metadata = {
        **dict(section.parser_metadata or {}),
        "layout_protocol_version": PARSER_LAYOUT_PROTOCOL_VERSION,
        "source_type": source_type,
        "layout_item_count": len(layout_items),
        "structure_object_count": len(structure_objects),
        "native_layout_available": any(bool(item.bbox) for item in layout_items),
    }
    return ParsedSection(
        title=title,
        text=text,
        page_number=section.page_number,
        section=section_label,
        metadata=metadata,
        layout_items=list(layout_items),
        structure_objects=list(structure_objects),
        parser_metadata=parser_metadata,
    )


def clean_parsed_sections(sections: list[ParsedSection], source_type: str, common_metadata: dict[str, Any] | None = None) -> list[ParsedSection]:
    return [
        cleaned
        for section in sections
        for cleaned in [_clean_section(section, source_type, common_metadata)]
        if cleaned.text
    ]


def _flow_parsed_section(
    *,
    title: str,
    text: str,
    section: str | None,
    page_number: int | None,
    content_kind: str,
    artifact_prefix: str,
    source: str,
    metadata: dict[str, Any] | None = None,
) -> ParsedSection:
    layouts, objects = _flow_artifacts(
        text,
        artifact_prefix=artifact_prefix,
        page_number=page_number,
        content_kind=content_kind,
        source=source,
        path_prefix=section or title,
    )
    return ParsedSection(
        title=title,
        text=text,
        page_number=page_number,
        section=section,
        metadata={"content_kind": content_kind, **(metadata or {})},
        layout_items=layouts,
        structure_objects=objects,
        parser_metadata={
            "parser": source,
            "layout_protocol_version": PARSER_LAYOUT_PROTOCOL_VERSION,
            "flow_block_protocol_version": FLOW_BLOCK_PROTOCOL_VERSION,
            "native_layout_available": False,
            **(
                {
                    "block_start_protocol_version": (
                        MARKDOWN_BLOCK_START_PROTOCOL_VERSION
                    ),
                    "html_block_protocol_version": (
                        MARKDOWN_HTML_BLOCK_PROTOCOL_VERSION
                    ),
                    "link_reference_protocol_version": (
                        MARKDOWN_LINK_REFERENCE_PROTOCOL_VERSION
                    ),
                    "target_spec_version": MARKDOWN_TARGET_SPEC_VERSION,
                }
                if content_kind == "markdown"
                else {}
            ),
        },
    )


def _markdown_section_text(buffer: list[str]) -> str:
    start = 0
    end = len(buffer)
    while start < end and not buffer[start].strip():
        start += 1
    while end > start and not buffer[end - 1].strip():
        end -= 1
    return "\n".join(buffer[start:end])


def parse_markdown(source: FrozenSourceSnapshot) -> list[ParsedSection]:
    source = _require_frozen_source(source)
    path = _parser_path(source)
    text, text_metadata = load_text_with_metadata(
        source,
        preserve_horizontal_whitespace=True,
        preserve_boundary_whitespace=True,
    )
    lines = _flow_line_spans(text)
    sections: list[ParsedSection] = []
    current_title = path.stem
    buffer: list[str] = []
    paragraph_open = False
    cursor = 0
    while cursor < len(lines):
        line = lines[cursor][2]
        opening_fence = _markdown_fence_marker(line)
        if opening_fence is not None:
            buffer.append(line)
            cursor += 1
            while cursor < len(lines):
                fence_line = lines[cursor][2]
                buffer.append(fence_line)
                cursor += 1
                if _markdown_fence_closes(fence_line, opening_fence):
                    break
            paragraph_open = False
            continue
        list_marker = _markdown_list_marker(line)
        if list_marker is not None and (
            not paragraph_open or _markdown_list_interrupts_paragraph(line)
        ):
            list_end = _markdown_list_block_end(lines, cursor)
            assert list_end is not None
            buffer.extend(
                lines[index][2]
                for index in range(cursor, list_end)
            )
            cursor = list_end
            paragraph_open = False
            continue
        if not paragraph_open:
            table_end = _markdown_table_end(lines, cursor)
            if table_end is not None:
                buffer.extend(
                    lines[index][2]
                    for index in range(cursor, table_end)
                )
                cursor = table_end
                paragraph_open = False
                continue
            link_reference_end = _markdown_link_reference_definition_end(
                lines,
                cursor,
            )
            if link_reference_end is not None:
                buffer.extend(
                    lines[index][2]
                    for index in range(cursor, link_reference_end)
                )
                cursor = link_reference_end
                paragraph_open = False
                continue
        html_block = _markdown_html_block_end(
            lines,
            cursor,
            paragraph_open=paragraph_open,
        )
        if html_block is not None:
            html_block_end, _html_block_type = html_block
            buffer.extend(
                lines[index][2]
                for index in range(cursor, html_block_end)
            )
            cursor = html_block_end
            paragraph_open = False
            continue
        heading_title = _markdown_atx_heading_title(line)
        if heading_title is not None:
            if buffer:
                section_text = _markdown_section_text(buffer)
                if section_text:
                    sections.append(
                        _flow_parsed_section(
                            title=current_title,
                            text=section_text,
                            section=current_title,
                            page_number=None,
                            content_kind="markdown",
                            artifact_prefix=f"markdown:{len(sections)}",
                            source="markdown_syntax_parser_v5",
                            metadata=text_metadata,
                        )
                    )
                buffer = []
            current_title = heading_title or current_title
            paragraph_open = False
            cursor += 1
            continue
        buffer.append(line)
        block_start_kind = _markdown_block_start_kind(line)
        if not line.strip():
            paragraph_open = False
        elif block_start_kind == "link_reference_definition":
            paragraph_open = paragraph_open
        elif block_start_kind == "indented_code":
            # GFM indented code cannot interrupt an open paragraph.
            paragraph_open = paragraph_open
        elif block_start_kind == "list":
            if not paragraph_open or _markdown_list_interrupts_paragraph(line):
                paragraph_open = False
        elif block_start_kind in {
            "block_quote",
            "thematic_break",
        }:
            paragraph_open = False
        else:
            paragraph_open = True
        cursor += 1
    if buffer:
        section_text = _markdown_section_text(buffer)
        if section_text:
            sections.append(
                _flow_parsed_section(
                    title=current_title,
                    text=section_text,
                    section=current_title,
                    page_number=None,
                    content_kind="markdown",
                    artifact_prefix=f"markdown:{len(sections)}",
                    source="markdown_syntax_parser_v5",
                    metadata=text_metadata,
                )
            )
    return [section for section in sections if section.text]


def parse_text(source: FrozenSourceSnapshot) -> list[ParsedSection]:
    source = _require_frozen_source(source)
    path = _parser_path(source)
    text, text_metadata = load_text_with_metadata(source)
    return [
        _flow_parsed_section(
            title=path.stem,
            text=text.strip(),
            section=path.stem,
            page_number=None,
            content_kind="text",
            artifact_prefix="text:0",
            source="plain_text_flow_parser_v1",
            metadata=text_metadata,
        )
    ]


def parse_html(source: FrozenSourceSnapshot) -> list[ParsedSection]:
    source = _require_frozen_source(source)
    path = _parser_path(source)
    html, text_metadata = load_text_with_metadata(source)
    soup = BeautifulSoup(html, "html.parser")
    title = soup.title.text.strip() if soup.title and soup.title.text else path.stem
    semantic_names = {"h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "pre", "table", "figcaption", "caption"}
    elements = [
        element
        for element in soup.find_all(semantic_names)
        if not any(getattr(parent, "name", None) in semantic_names for parent in element.parents if parent is not soup)
    ]
    parts: list[str] = []
    layouts: list[ParsedLayoutItem] = []
    objects: list[ParsedStructureObject] = []
    cursor = 0
    heading_path: list[str] = []
    for order, element in enumerate(elements):
        block_text = element.get_text("\n" if element.name in {"pre", "table"} else " ", strip=True)
        if not block_text:
            continue
        if parts:
            parts.append("\n\n")
            cursor += 2
        start = cursor
        parts.append(block_text)
        cursor += len(block_text)
        end = cursor
        if element.name and element.name.startswith("h") and len(element.name) == 2:
            level = max(1, min(6, int(element.name[1])))
            heading_path = heading_path[: level - 1] + [block_text]
            object_type = "section"
        else:
            object_type = {
                "li": "list",
                "pre": "code_block",
                "table": "table",
                "figcaption": "caption",
                "caption": "caption",
            }.get(element.name, "paragraph")
        metadata = {
            "parser_source": "beautifulsoup_dom_v1",
            "html_tag": element.name,
            "native_geometry": False,
            "layout_protocol_version": PARSER_LAYOUT_PROTOCOL_VERSION,
        }
        layout_id = f"html:layout:{order}"
        layouts.append(
            ParsedLayoutItem(
                layout_id=layout_id,
                text=block_text,
                char_start=start,
                char_end=end,
                coordinate_system="html_dom_order_v1",
                region_type=object_type,
                reading_order=order,
                metadata=metadata,
            )
        )
        objects.append(
            ParsedStructureObject(
                structure_id=f"html:structure:{order}",
                object_type=object_type,
                text=block_text,
                char_start=start,
                char_end=end,
                title=block_text[:120],
                coordinate_system="html_dom_order_v1",
                reading_order=order,
                path=" / ".join(heading_path) or title,
                metadata=metadata,
            )
        )
    text = "".join(parts) or soup.get_text("\n", strip=True)
    if not layouts:
        layouts, objects = _flow_artifacts(
            text,
            artifact_prefix="html",
            page_number=None,
            content_kind="html",
            coordinate_system="html_dom_order_v1",
            source="beautifulsoup_dom_v1",
        )
    return [
        ParsedSection(
            title=title,
            text=text,
            section=title,
            metadata={"content_kind": "html", **text_metadata},
            layout_items=layouts,
            structure_objects=objects,
            parser_metadata={
                "parser": "beautifulsoup_dom_v1",
                "layout_protocol_version": PARSER_LAYOUT_PROTOCOL_VERSION,
                "native_layout_available": False,
            },
        )
    ]


def _detect_formula(text: str) -> bool:
    """Detect formula-heavy text without rewriting math symbols."""
    formula_chars = set("∑∫?√∞≈≠≤≥±×÷∈???∪∩→←????αβγδθηλμπρστυφχψωΓΔΘΛΠΣΦΨΩ")
    if not text:
        return False
    ratio = sum(1 for c in text if c in formula_chars) / len(text)
    return ratio > 0.03


def ocr_image_to_text(image: Any, *, source_name: str) -> tuple[str, dict[str, Any]]:
    """Run OCR on an image path or PIL image and return text plus audit metadata."""
    from PIL import Image

    pil_image = Image.open(image) if isinstance(image, Path) else image
    image_width, image_height = pil_image.size
    if isinstance(image, Path):
        try:
            from paddleocr import PaddleOCR

            ocr = PaddleOCR(use_angle_cls=True, lang="ch")
            results = ocr.ocr(str(image), cls=True)
            rows: list[dict[str, Any]] = []
            for page in results:
                for line in page:
                    points = line[0]
                    line_text = str(line[1][0]).strip()
                    confidence = float(line[1][1] or 0.0)
                    if not line_text or not points:
                        continue
                    xs = [float(point[0]) for point in points]
                    ys = [float(point[1]) for point in points]
                    rows.append(
                        {
                            "text": line_text,
                            "bbox": _normalized_native_bbox(
                                min(xs),
                                min(ys),
                                max(xs),
                                max(ys),
                                width=image_width,
                                height=image_height,
                                page_number=1,
                                raw_coordinate_system="image_pixels_top_left",
                            ),
                            "confidence": confidence,
                        }
                    )
            text = "\n".join(row["text"] for row in rows)
            if text.strip():
                return text.strip(), {
                    "ocr_engine": "paddleocr",
                    "ocr_confidence": round(sum(row["confidence"] for row in rows) / max(len(rows), 1), 6),
                    "ocr_layout_items": rows,
                    "image_size": [image_width, image_height],
                }
        except Exception:
            pass

    try:
        import pytesseract
        from pytesseract import Output

        data = pytesseract.image_to_data(pil_image, lang="chi_sim+eng", output_type=Output.DICT)
        grouped: dict[tuple[int, int, int], list[int]] = {}
        for index, token in enumerate(data.get("text") or []):
            if not str(token).strip():
                continue
            key = (
                int((data.get("block_num") or [0])[index]),
                int((data.get("par_num") or [0])[index]),
                int((data.get("line_num") or [0])[index]),
            )
            grouped.setdefault(key, []).append(index)
        rows = []
        for indices in grouped.values():
            line_text = " ".join(str(data["text"][index]).strip() for index in indices).strip()
            left = min(int(data["left"][index]) for index in indices)
            top = min(int(data["top"][index]) for index in indices)
            right = max(int(data["left"][index]) + int(data["width"][index]) for index in indices)
            bottom = max(int(data["top"][index]) + int(data["height"][index]) for index in indices)
            confidences = [float(data["conf"][index]) for index in indices if float(data["conf"][index]) >= 0]
            rows.append(
                {
                    "text": line_text,
                    "bbox": _normalized_native_bbox(
                        left,
                        top,
                        right,
                        bottom,
                        width=image_width,
                        height=image_height,
                        page_number=1,
                        raw_coordinate_system="image_pixels_top_left",
                    ),
                    "confidence": (sum(confidences) / max(len(confidences), 1)) / 100.0,
                }
            )
        text = "\n".join(row["text"] for row in rows)
        if text.strip():
            return text.strip(), {
                "ocr_engine": "pytesseract",
                "ocr_confidence": round(sum(row["confidence"] for row in rows) / max(len(rows), 1), 6),
                "ocr_layout_items": rows,
                "image_size": [image_width, image_height],
            }
    except Exception as exc:
        raise RuntimeError(
            f"OCR dependencies unavailable for {source_name}: install/enable the api OCR extra "
            "or provide a working PaddleOCR/Tesseract runtime"
        ) from exc
    raise RuntimeError(f"No OCR text extracted from {source_name}")


def _normalized_native_bbox(
    x0: float,
    y0: float,
    x1: float,
    y1: float,
    *,
    width: float,
    height: float,
    page_number: int,
    raw_coordinate_system: str,
) -> dict[str, Any]:
    safe_width = max(float(width), 1e-9)
    safe_height = max(float(height), 1e-9)
    return {
        "page_number": page_number,
        "x0": round(max(0.0, min(1.0, float(x0) / safe_width)), 6),
        "y0": round(max(0.0, min(1.0, float(y0) / safe_height)), 6),
        "x1": round(max(0.0, min(1.0, float(x1) / safe_width)), 6),
        "y1": round(max(0.0, min(1.0, float(y1) / safe_height)), 6),
        "coordinate_system": "normalized_page_v1",
        "synthetic": False,
        "raw_bbox": [float(x0), float(y0), float(x1), float(y1)],
        "raw_coordinate_system": raw_coordinate_system,
        "page_size": [float(width), float(height)],
    }


def parse_pdf(source: FrozenSourceSnapshot) -> list[ParsedSection]:
    import fitz
    from PIL import Image

    source = _require_frozen_source(source)
    path = _parser_path(source)
    sections: list[ParsedSection] = []
    with fitz.open(stream=source.content_bytes, filetype="pdf") as document:
        for idx, page in enumerate(document, start=1):
            page_width = float(page.rect.width)
            page_height = float(page.rect.height)
            native_blocks = []
            try:
                native_blocks = [block for block in page.get_text("blocks", sort=True) if len(block) >= 5 and str(block[4]).strip()]
            except Exception:
                native_blocks = []
            image_list = page.get_images(full=True)
            image_errors: list[str] = []
            block_rows: list[dict[str, Any]] = []
            for block_index, block in enumerate(native_blocks):
                block_text = str(block[4]).strip()
                if not block_text:
                    continue
                block_type = int(block[6]) if len(block) > 6 and str(block[6]).isdigit() else 0
                if block_type != 0:
                    continue
                bbox = _normalized_native_bbox(
                    float(block[0]),
                    float(block[1]),
                    float(block[2]),
                    float(block[3]),
                    width=page_width,
                    height=page_height,
                    page_number=idx,
                    raw_coordinate_system="pdf_points_top_left",
                )
                block_rows.append(
                    {
                        "text": block_text,
                        "bbox": bbox,
                        "source": "pymupdf_text_block",
                        "source_index": block_index,
                        "confidence": 1.0,
                    }
                )
            for img_idx, img in enumerate(image_list, start=1):
                xref = img[0]
                try:
                    base_image = document.extract_image(xref)
                    image = Image.open(io.BytesIO(base_image["image"]))
                    img_ocr_text, img_ocr_meta = ocr_image_to_text(
                        image,
                        source_name=f"{path.name} p.{idx} img.{img_idx}",
                    )
                    if img_ocr_text:
                        rects = page.get_image_rects(xref)
                        rect = rects[0] if rects else page.rect
                        block_rows.append(
                            {
                                "text": img_ocr_text,
                                "bbox": _normalized_native_bbox(
                                    rect.x0,
                                    rect.y0,
                                    rect.x1,
                                    rect.y1,
                                    width=page_width,
                                    height=page_height,
                                    page_number=idx,
                                    raw_coordinate_system="pdf_points_top_left",
                                ),
                                "source": "embedded_image_ocr",
                                "source_index": img_idx,
                                "confidence": float(img_ocr_meta.get("ocr_confidence") or 0.8),
                                "ocr_metadata": img_ocr_meta,
                            }
                        )
                except Exception as exc:
                    image_errors.append(str(exc))
                    continue
            if not block_rows:
                continue

            parts: list[str] = []
            layouts: list[ParsedLayoutItem] = []
            objects: list[ParsedStructureObject] = []
            cursor = 0
            for reading_order, row in enumerate(block_rows):
                if parts:
                    parts.append("\n\n")
                    cursor += 2
                block_text = str(row["text"]).strip()
                start = cursor
                parts.append(block_text)
                cursor += len(block_text)
                end = cursor
                object_type = _block_object_type(block_text, content_kind="ocr" if row["source"] == "embedded_image_ocr" else "pdf_page")
                artifact_metadata = {
                    "parser_source": row["source"],
                    "source_index": row["source_index"],
                    "native_geometry": True,
                    "layout_protocol_version": PARSER_LAYOUT_PROTOCOL_VERSION,
                    **dict(row.get("ocr_metadata") or {}),
                }
                layouts.append(
                    ParsedLayoutItem(
                        layout_id=f"pdf:p{idx}:layout:{reading_order}",
                        text=block_text,
                        char_start=start,
                        char_end=end,
                        page_number=idx,
                        bbox=dict(row["bbox"]),
                        coordinate_system="normalized_page_v1",
                        region_type=object_type,
                        reading_order=reading_order,
                        confidence=float(row["confidence"]),
                        metadata=artifact_metadata,
                    )
                )
                objects.append(
                    ParsedStructureObject(
                        structure_id=f"pdf:p{idx}:structure:{reading_order}",
                        object_type=object_type,
                        text=block_text,
                        char_start=start,
                        char_end=end,
                        title=re.sub(r"\s+", " ", block_text).strip()[:120] or object_type,
                        page_number=idx,
                        bbox=dict(row["bbox"]),
                        coordinate_system="normalized_page_v1",
                        reading_order=reading_order,
                        confidence=float(row["confidence"]),
                        metadata=artifact_metadata,
                    )
                )
            text = "".join(parts)
            lines = text.splitlines()
            has_table = any("|" in line and "---" in line for line in lines)
            has_formula = _detect_formula(text)
            page_title = lines[0][:120] if lines else ""
            metadata: dict[str, Any] = {
                "content_kind": "pdf_page",
                "pdf_image_count": len(image_list),
                "native_layout_block_count": len(layouts),
                "page_size": [page_width, page_height],
            }
            if has_table:
                metadata["has_table"] = True
            if has_formula:
                metadata["has_formula"] = True
            if any(row["source"] == "embedded_image_ocr" for row in block_rows):
                metadata.update({"ocr_applied": True, "ocr_page_count": 1, "ocr_reason": "page_contains_images"})
            if image_errors:
                metadata["ocr_image_errors"] = image_errors

            sections.append(
                ParsedSection(
                    title=page_title or f"{path.stem} p.{idx}",
                    text=text,
                    page_number=idx,
                    section=page_title or path.stem,
                    metadata=metadata,
                    layout_items=layouts,
                    structure_objects=objects,
                    parser_metadata={
                        "parser": "pymupdf_blocks_with_ocr_v1",
                        "layout_protocol_version": PARSER_LAYOUT_PROTOCOL_VERSION,
                        "native_layout_available": True,
                        "page_size": [page_width, page_height],
                    },
                )
            )
    return sections


def parse_presentation(source: FrozenSourceSnapshot) -> list[ParsedSection]:
    from pptx import Presentation

    source = _require_frozen_source(source)
    path = _parser_path(source)
    presentation = Presentation(io.BytesIO(source.content_bytes))
    sections: list[ParsedSection] = []
    slide_width = float(presentation.slide_width)
    slide_height = float(presentation.slide_height)
    for idx, slide in enumerate(presentation.slides, start=1):
        shape_rows: list[dict[str, Any]] = []
        for source_index, shape in enumerate(slide.shapes):
            shape_text = str(getattr(shape, "text", "") or "").strip()
            if bool(getattr(shape, "has_table", False)):
                table_rows = [
                    " | ".join(str(cell.text or "").strip() for cell in row.cells)
                    for row in shape.table.rows
                ]
                shape_text = "\n".join(row for row in table_rows if row.strip()) or shape_text
            if not shape_text:
                continue
            placeholder_type = ""
            if bool(getattr(shape, "is_placeholder", False)):
                try:
                    placeholder_type = str(shape.placeholder_format.type)
                except Exception:
                    placeholder_type = ""
            object_type = (
                "table"
                if bool(getattr(shape, "has_table", False))
                else "section"
                if "TITLE" in placeholder_type.upper()
                else _block_object_type(shape_text, content_kind="slide")
            )
            shape_rows.append(
                {
                    "source_index": source_index,
                    "text": shape_text,
                    "object_type": object_type,
                    "placeholder_type": placeholder_type,
                    "left": float(getattr(shape, "left", 0) or 0),
                    "top": float(getattr(shape, "top", 0) or 0),
                    "width": float(getattr(shape, "width", slide_width) or slide_width),
                    "height": float(getattr(shape, "height", slide_height) or slide_height),
                }
            )
        shape_rows.sort(key=lambda row: (row["top"], row["left"], row["source_index"]))
        parts: list[str] = []
        layouts: list[ParsedLayoutItem] = []
        objects: list[ParsedStructureObject] = []
        cursor = 0
        for reading_order, row in enumerate(shape_rows):
            if parts:
                parts.append("\n\n")
                cursor += 2
            start = cursor
            parts.append(row["text"])
            cursor += len(row["text"])
            end = cursor
            bbox = _normalized_native_bbox(
                row["left"],
                row["top"],
                row["left"] + row["width"],
                row["top"] + row["height"],
                width=slide_width,
                height=slide_height,
                page_number=idx,
                raw_coordinate_system="pptx_emu_top_left",
            )
            metadata = {
                "parser_source": "python_pptx_shape_v1",
                "source_shape_index": row["source_index"],
                "placeholder_type": row["placeholder_type"],
                "native_geometry": True,
                "layout_protocol_version": PARSER_LAYOUT_PROTOCOL_VERSION,
            }
            layouts.append(
                ParsedLayoutItem(
                    layout_id=f"pptx:s{idx}:layout:{reading_order}",
                    text=row["text"],
                    char_start=start,
                    char_end=end,
                    page_number=idx,
                    bbox=bbox,
                    coordinate_system="normalized_page_v1",
                    region_type=row["object_type"],
                    reading_order=reading_order,
                    metadata=metadata,
                )
            )
            objects.append(
                ParsedStructureObject(
                    structure_id=f"pptx:s{idx}:structure:{reading_order}",
                    object_type=row["object_type"],
                    text=row["text"],
                    char_start=start,
                    char_end=end,
                    title=re.sub(r"\s+", " ", row["text"]).strip()[:120],
                    page_number=idx,
                    bbox=bbox,
                    coordinate_system="normalized_page_v1",
                    reading_order=reading_order,
                    metadata=metadata,
                )
            )
        text = "".join(parts)
        if text:
            lines = text.splitlines()
            sections.append(
                ParsedSection(
                    title=lines[0][:120] if lines else f"{path.stem} slide {idx}",
                    text=text,
                    page_number=idx,
                    section=f"slide-{idx}",
                    metadata={"content_kind": "slide", "native_layout_shape_count": len(layouts)},
                    layout_items=layouts,
                    structure_objects=objects,
                    parser_metadata={
                        "parser": "python_pptx_shape_v1",
                        "layout_protocol_version": PARSER_LAYOUT_PROTOCOL_VERSION,
                        "native_layout_available": True,
                        "page_size": [slide_width, slide_height],
                    },
                )
            )
    return sections


def parse_docx(source: FrozenSourceSnapshot) -> list[ParsedSection]:
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    source = _require_frozen_source(source)
    path = _parser_path(source)
    document = Document(io.BytesIO(source.content_bytes))
    sections: list[ParsedSection] = []
    current_title = path.stem
    blocks: list[dict[str, Any]] = []

    def flush_section() -> None:
        nonlocal blocks
        if not blocks:
            return
        parts: list[str] = []
        layouts: list[ParsedLayoutItem] = []
        objects: list[ParsedStructureObject] = []
        cursor = 0
        for reading_order, block in enumerate(blocks):
            if parts:
                parts.append("\n\n")
                cursor += 2
            start = cursor
            parts.append(block["text"])
            cursor += len(block["text"])
            end = cursor
            metadata = {
                "parser_source": "python_docx_flow_v1",
                "source_block_index": block["source_index"],
                "style": block.get("style"),
                "native_geometry": False,
                "layout_protocol_version": PARSER_LAYOUT_PROTOCOL_VERSION,
            }
            layouts.append(
                ParsedLayoutItem(
                    layout_id=f"docx:section:{len(sections)}:layout:{reading_order}",
                    text=block["text"],
                    char_start=start,
                    char_end=end,
                    coordinate_system="docx_flow_v1",
                    region_type=block["object_type"],
                    reading_order=reading_order,
                    metadata=metadata,
                )
            )
            objects.append(
                ParsedStructureObject(
                    structure_id=f"docx:section:{len(sections)}:structure:{reading_order}",
                    object_type=block["object_type"],
                    text=block["text"],
                    char_start=start,
                    char_end=end,
                    title=re.sub(r"\s+", " ", block["text"]).strip()[:120],
                    coordinate_system="docx_flow_v1",
                    reading_order=reading_order,
                    path=current_title,
                    metadata=metadata,
                )
            )
        section_text = "".join(parts)
        sections.append(
            ParsedSection(
                title=current_title,
                text=section_text,
                section=current_title,
                metadata={"content_kind": "doc_section", "native_structure_block_count": len(objects)},
                layout_items=layouts,
                structure_objects=objects,
                parser_metadata={
                    "parser": "python_docx_flow_v1",
                    "layout_protocol_version": PARSER_LAYOUT_PROTOCOL_VERSION,
                    "native_layout_available": False,
                },
            )
        )
        blocks = []

    for source_index, child in enumerate(document.element.body.iterchildren()):
        if child.tag.endswith("}p"):
            paragraph = Paragraph(child, document)
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = str(paragraph.style.name if paragraph.style else "")
            if "Heading" in style_name:
                flush_section()
                current_title = text
                continue
            object_type = "list" if style_name.lower().startswith("list") else "paragraph"
            blocks.append(
                {
                    "text": text,
                    "object_type": object_type,
                    "source_index": source_index,
                    "style": style_name,
                }
            )
        elif child.tag.endswith("}tbl"):
            table = Table(child, document)
            table_text = "\n".join(
                " | ".join(str(cell.text or "").strip() for cell in row.cells)
                for row in table.rows
            ).strip()
            if table_text:
                blocks.append(
                    {
                        "text": table_text,
                        "object_type": "table",
                        "source_index": source_index,
                        "style": "table",
                    }
                )
    flush_section()
    return sections


def parse_notebook(source: FrozenSourceSnapshot) -> list[ParsedSection]:
    source = _require_frozen_source(source)
    path = _parser_path(source)
    decoded, text_metadata = decode_text_bytes(source.content_bytes)
    notebook = json.loads(decoded)
    sections: list[ParsedSection] = []
    current_title = path.stem
    for current_index, cell in enumerate(notebook.get("cells", []), start=1):
        source = "".join(cell.get("source", []))
        if not source.strip():
            continue
        cell_type = str(cell.get("cell_type") or "raw")
        if cell_type == "markdown":
            heading = next((line.lstrip("# ").strip() for line in source.splitlines() if line.strip().startswith("#")), None)
            if heading:
                current_title = heading
            cell_text = source.strip()
            content_kind = "markdown"
            title = current_title
        else:
            outputs: list[str] = []
            for output in cell.get("outputs", []):
                if "text" in output:
                    outputs.append("".join(output["text"]))
                elif "data" in output and "text/plain" in output["data"]:
                    outputs.append("".join(output["data"]["text/plain"]))
            block = ["[Code Cell]", source.strip()]
            if outputs:
                block.extend(["[Output]", "\n".join(outputs).strip()])
            cell_text = "\n".join(part for part in block if part)
            content_kind = "code"
            title = f"{current_title} code"
        layouts, objects = _flow_artifacts(
            cell_text,
            artifact_prefix=f"notebook:cell:{current_index}",
            page_number=None,
            content_kind=content_kind,
            coordinate_system="notebook_cell_order_v1",
            source="nbformat_cell_parser_v1",
        )
        sections.append(
            ParsedSection(
                title=title,
                text=cell_text,
                section=current_title,
                metadata={"cell_index": current_index, "cell_type": cell_type, "content_kind": content_kind, **text_metadata},
                layout_items=layouts,
                structure_objects=objects,
                parser_metadata={
                    "parser": "nbformat_cell_parser_v1",
                    "layout_protocol_version": PARSER_LAYOUT_PROTOCOL_VERSION,
                    "native_layout_available": False,
                    "cell_index": current_index,
                },
            )
        )
    return sections


def parse_image(source: FrozenSourceSnapshot) -> list[ParsedSection]:
    from PIL import Image

    source = _require_frozen_source(source)
    path = _parser_path(source)
    image = Image.open(io.BytesIO(source.content_bytes))
    image.load()
    text, metadata = ocr_image_to_text(image, source_name=path.name)
    layouts: list[ParsedLayoutItem] = []
    objects: list[ParsedStructureObject] = []
    cursor = 0
    for reading_order, row in enumerate(metadata.get("ocr_layout_items") or []):
        row_text = str(row.get("text") or "").strip()
        start = text.find(row_text, cursor)
        if not row_text or start < 0:
            continue
        end = start + len(row_text)
        artifact_metadata = {
            "parser_source": str(metadata.get("ocr_engine") or "ocr"),
            "native_geometry": True,
            "layout_protocol_version": PARSER_LAYOUT_PROTOCOL_VERSION,
        }
        layouts.append(
            ParsedLayoutItem(
                layout_id=f"image:layout:{reading_order}",
                text=row_text,
                char_start=start,
                char_end=end,
                page_number=1,
                bbox=dict(row.get("bbox") or {}),
                coordinate_system="normalized_page_v1",
                region_type="ocr_region",
                reading_order=reading_order,
                confidence=float(row.get("confidence") or metadata.get("ocr_confidence") or 0.8),
                metadata=artifact_metadata,
            )
        )
        objects.append(
            ParsedStructureObject(
                structure_id=f"image:structure:{reading_order}",
                object_type=_block_object_type(row_text, content_kind="ocr"),
                text=row_text,
                char_start=start,
                char_end=end,
                title=row_text[:120],
                page_number=1,
                bbox=dict(row.get("bbox") or {}),
                coordinate_system="normalized_page_v1",
                reading_order=reading_order,
                confidence=float(row.get("confidence") or metadata.get("ocr_confidence") or 0.8),
                metadata=artifact_metadata,
            )
        )
        cursor = end
    if not layouts:
        full_bbox = {
            "page_number": 1,
            "x0": 0.0,
            "y0": 0.0,
            "x1": 1.0,
            "y1": 1.0,
            "coordinate_system": "normalized_page_v1",
            "synthetic": False,
            "raw_coordinate_system": "image_pixels_top_left",
            "page_size": metadata.get("image_size") or [],
        }
        layouts, objects = _flow_artifacts(
            text,
            artifact_prefix="image",
            page_number=1,
            content_kind="ocr",
            coordinate_system="normalized_page_v1",
            bbox=full_bbox,
            source=str(metadata.get("ocr_engine") or "ocr"),
        )
    return [
        ParsedSection(
            title=path.stem,
            text=text,
            page_number=1,
            section=path.stem,
            metadata={"content_kind": "ocr", "ocr_applied": True, "ocr_page_count": 1, "ocr_reason": "image_file", **metadata},
            layout_items=layouts,
            structure_objects=objects,
            parser_metadata={
                "parser": str(metadata.get("ocr_engine") or "ocr"),
                "layout_protocol_version": PARSER_LAYOUT_PROTOCOL_VERSION,
                "native_layout_available": True,
                "image_size": metadata.get("image_size") or [],
            },
        )
    ]


def parse_with_unstructured(source: FrozenSourceSnapshot) -> list[ParsedSection]:
    from unstructured.partition.auto import partition

    source = _require_frozen_source(source)
    path = _parser_path(source)
    elements = partition(
        file=io.BytesIO(source.content_bytes),
        metadata_filename=path.name,
    )
    text = "\n".join(str(element) for element in elements if str(element).strip())
    return [
        _flow_parsed_section(
            title=path.stem,
            text=text,
            section=path.stem,
            page_number=None,
            content_kind="unstructured",
            artifact_prefix="unstructured:0",
            source="unstructured_partition_v1",
        )
    ] if text else []


def parse_document(source: FrozenSourceSnapshot) -> tuple[str, list[ParsedSection]]:
    source = _require_frozen_source(source)
    path = _parser_path(source)
    source_type = detect_source_type(source)
    parsers = {
        "pdf": parse_pdf,
        "ppt": parse_presentation,
        "pptx": parse_presentation,
        "docx": parse_docx,
        "markdown": parse_markdown,
        "text": parse_text,
        "image": parse_image,
        "notebook": parse_notebook,
        "html": parse_html,
    }
    parser = parsers.get(source_type)
    if parser is None:
        if get_settings().enable_model_fallback:
            return source_type, clean_parsed_sections(parse_with_unstructured(source), source_type)
        raise RuntimeError(f"Unsupported source type for {path.name}: {source_type}")
    try:
        return source_type, clean_parsed_sections(parser(source), source_type)
    except Exception as exc:
        if get_settings().enable_model_fallback:
            return source_type, clean_parsed_sections(parse_with_unstructured(source), source_type)
        raise RuntimeError(f"Failed to parse {path.name} as {source_type}: {exc}") from exc


def sections_to_json(sections: list[ParsedSection]) -> list[dict[str, Any]]:
    return [asdict(section) for section in sections]


INVALID_PARTITION_LABELS = {
    "data",
    "storage",
    "reviewmarkdown",
}


def _normalize_label(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", " ", value.lower()).strip()


def is_date_like_label(value: str) -> bool:
    return bool(re.fullmatch(r"(?:19|20)\d{6}", value.strip()))


def is_invalid_partition_label(value: str | None, knowledge_base_name: str | None = None) -> bool:
    if not value:
        return True
    normalized = _normalize_label(value)
    if not normalized:
        return True
    if normalized.replace(" ", "") in INVALID_PARTITION_LABELS:
        return True
    if is_date_like_label(normalized.replace(" ", "")):
        return True
    if knowledge_base_name and normalized == _normalize_label(knowledge_base_name):
        return True
    return False


def canonical_partition_label(value: str, knowledge_base_name: str | None = None) -> str | None:
    cleaned = re.sub(r"[_-]+", " ", value)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    if not cleaned:
        return None

    match = re.search(r"\b(?:partition|chap)\s*\.?\s*(\d+[A-Za-z]?)\b", cleaned, flags=re.IGNORECASE)
    if match:
        return f"partition {match.group(1)}"

    match = re.search(r"\b(?:lecture|lec|l)\s*\.?\s*(\d+[A-Za-z]?)\b", cleaned, flags=re.IGNORECASE)
    if match:
        return f"Lecture {match.group(1)}"

    match = re.search(r"\bweek\s*\.?\s*(\d+[A-Za-z]?)\b", cleaned, flags=re.IGNORECASE)
    if match:
        return f"Week {match.group(1)}"

    match = re.search(r"\blab\s*\.?\s*(\d+[A-Za-z]?)\b", cleaned, flags=re.IGNORECASE)
    if match:
        return f"Lab {match.group(1)}"

    if re.search(r"\blabs?\b|\blaboratory\b", cleaned, flags=re.IGNORECASE):
        if re.search(r"\bquestions?\b", cleaned, flags=re.IGNORECASE):
            return "Lab Questions"
        if re.search(r"\bsolutions?\b", cleaned, flags=re.IGNORECASE):
            return "Lab Solutions"
        return "Lab"

    if re.search(r"\bwork\s*book\b|\bworkbook\b|\bwork\s*items?\b", cleaned, flags=re.IGNORECASE):
        return "Workbook"

    if re.search(r"\bz\s*table\b|\breference\b|\bformula\b|\bsummary\b|\bvisuali[sz]er\b", cleaned, flags=re.IGNORECASE):
        return "Reference"

    if re.search(r"\breview\b|\brevision\b|\brecap\b", cleaned, flags=re.IGNORECASE):
        return "Review"

    cleaned = cleaned[:80].strip()
    if cleaned and not re.search(r"[A-Za-z0-9]", cleaned):
        return "Reference"
    return None if is_invalid_partition_label(cleaned, knowledge_base_name=knowledge_base_name) else cleaned


def derive_partition(path: Path, knowledge_base_name: str | None = None) -> str:
    candidates = [path.stem, path.parent.name]
    for item in candidates:
        label = canonical_partition_label(item, knowledge_base_name=knowledge_base_name)
        if label and not is_invalid_partition_label(label, knowledge_base_name=knowledge_base_name):
            return label
    return path.stem[:80]
